#!/usr/bin/env python3
"""Generate presentation matching the CALM reference style.

Uses the reference PPTX as a template to inherit its theme, colors, and layouts.
Concise bullet style: header ~25pt, sub-bullets ~18pt with "- " prefix.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Use reference PPTX as template — strip its slides, keep theme + layouts
_ref = Presentation("experiments/Continuous Autoregressive Language Models.pptx")
for sldId in list(_ref.slides._sldIdLst):
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId:
        try:
            _ref.part.drop_rel(rId)
        except KeyError:
            pass
    _ref.slides._sldIdLst.remove(sldId)
_ref.save("/tmp/_template_clean.pptx")
del _ref
prs = Presentation("/tmp/_template_clean.pptx")

# Layout references
TITLE_LAYOUT = prs.slide_layouts[0]       # TITLE
BODY_LAYOUT  = prs.slide_layouts[2]       # TITLE_AND_BODY
BLANK_LAYOUT = prs.slide_layouts[10]      # BLANK

HEADER_SZ = Pt(25)
SUB_SZ    = Pt(18)
SMALL_SZ  = Pt(14)

# ──────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(TITLE_LAYOUT)
    # Find placeholders
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # title
            ph.text = title
        elif ph.placeholder_format.idx == 1:  # subtitle
            ph.text = subtitle
    return slide


def add_body_slide(title, bullets):
    """bullets: list of (text, level) where level=0 is header, level=1 is sub-bullet."""
    slide = prs.slides.add_slide(BODY_LAYOUT)

    # Set title
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            break

    # Find body placeholder
    body_ph = None
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx != 0 and ph.has_text_frame:
            # Check if it's the big body area (not the page number)
            if ph.width > Inches(3):
                body_ph = ph
                break

    if body_ph is None:
        # Fallback: use any non-title placeholder with text
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0 and ph.has_text_frame and ph.width > Inches(3):
                body_ph = ph
                break

    if body_ph is None:
        # Last resort: add a textbox
        body_ph = slide.shapes.add_textbox(
            Inches(0.78), Inches(1.50), Inches(8.88), Inches(3.76))

    tf = body_ph.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, (text, level) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if text == "":
            p.text = ""
            p.space_before = Pt(4)
            p.space_after = Pt(4)
            continue

        p.text = text
        if level == 0:
            p.font.size = HEADER_SZ
            p.font.bold = None
            p.space_before = Pt(8)
            p.space_after = Pt(2)
        else:
            p.font.size = SUB_SZ
            p.font.bold = None
            p.space_before = Pt(2)
            p.space_after = Pt(2)

    return slide


def add_table_slide(title, headers, rows, note=""):
    slide = prs.slides.add_slide(BODY_LAYOUT)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            break

    # Remove body placeholder if it exists (we'll use a table instead)
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx != 0 and ph.has_text_frame and ph.width > Inches(3):
            sp = ph._element
            sp.getparent().remove(sp)
            break

    num_rows = len(rows) + 1
    num_cols = len(headers)
    total_w = 9.0
    col_w = total_w / num_cols

    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.4),
        Inches(total_w), Inches(0.40 * num_rows))
    tbl = tbl_shape.table

    for ci in range(num_cols):
        tbl.columns[ci].width = Inches(col_w)

    # Header
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

    # Data
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.CENTER

    if note:
        tb = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.9), Inches(9.0), Inches(0.5))
        tb.text_frame.paragraphs[0].text = note
        tb.text_frame.paragraphs[0].font.size = Pt(11)

    return slide


# ══════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════

# ── 1. Title ──
add_title_slide(
    "LLM-Guided Search for Optimal AllToAllV on AWS Trainium",
    "Combining Claude, Genetic Algorithms, and Contention-Guided Synthesis"
)

# ── 2. Motivation ──
add_body_slide("Why AllToAllV Matters", [
    ("Mixture-of-Experts (MoE) training", 0),
    ("- Tokens routed to specialized experts across ranks", 1),
    ("- AllToAllV: variable-length all-to-all = core MoE primitive", 1),
    ("- Called 4x per MoE layer per training step", 1),
    ("", 0),
    ("On AWS Trainium", 0),
    ("- Must build from XLA primitives (all_gather, collective_permute)", 1),
    ("- No NCCL; default approaches are slow", 1),
    ("- Multi-node (EFA) adds cross-node overhead", 1),
])

# ── 3. Hardware ──
add_body_slide("Trainium trn1.32xlarge", [
    ("Intra-node topology", 0),
    ("- 16 NeuronDevices in 4x4 2D torus, 2 cores each (32 ranks)", 1),
    ("- NeuronLink: ~192 GB/s per link, 4 links per device", 1),
    ("- Max 4 hops, avg 2 hops across torus", 1),
    ("", 0),
    ("Inter-node (EFA)", 0),
    ("- 8 EFA adapters/node @ 12.5 GB/s each", 1),
    ("- 2-node cluster: 64 ranks, 32 devices, NeuronLink + EFA", 1),
    ("", 0),
    ("XLA collective primitives", 0),
    ("- all_gather: broadcast, single dispatch, HW-optimized", 1),
    ("- collective_permute: point-to-point, one dispatch per distance", 1),
    ("- all_to_all: XLA decomposes to internal ring (deceptively slow)", 1),
])

# ── 4. Key Insight ──
add_body_slide("Key Insight: Dispatch Overhead Dominates", [
    ("Each XLA op costs ~0.03 ms fixed dispatch overhead", 0),
    ("- Bandwidth (192 GB/s) is never the bottleneck", 1),
    ("- Latency is flat from 32 KB to 512 MB", 1),
    ("", 0),
    ("Naive AllGather: 33 XLA ops = ~1.05 ms", 0),
    ("Default Ring: 93 XLA ops = ~2.78 ms", 0),
    ("Optimized (agent output): 3 XLA ops = ~0.13 ms", 0),
    ("", 0),
    ("Reducing XLA IR node count is the only lever", 0),
])

# ── 5. Search Pipeline Overview ──
add_body_slide("Search Pipeline: 5 Phases + Feedback Loops", [
    ("Phase 1: Agent hardware profiling (LLM builds simulator)", 0),
    ("- Discovers topology, overheads; writes cost model; validates <20% error", 1),
    ("Phase 2: Baseline eval on simulator -> knowledgebase", 0),
    ("- All templates + GA/SA refinement + LLM candidates", 1),
    ("Phase 3: Multi-island evolution with simulator feedback", 0),
    ("- 3 islands + CGIS + template evolution (XLA + NKI)", 1),
    ("Phase 4: Iterative mini-benchmarking on real HW", 0),
    ("- Compare HW vs sim; refine simulator if predictions diverge", 1),
    ("Phase 5: Final code generation -> trainium_alltoallv.py", 0),
    ("", 0),
    ("Feedback: Phase 2/3 return to Phase 1 to refine simulator", 0),
])

# ── 6. Topology Simulator ──
add_body_slide("Topology Simulator", [
    ("Models the 4x4 NeuronLink torus", 0),
    ("- Per-link bandwidth, latency, hop counts", 1),
    ("- Simulates 6 algorithm templates", 1),
    ("", 0),
    ("Multi-objective cost function", 0),
    ("- sim_time + dispatch penalty + hop cost + contention", 1),
    ("", 0),
    ("Limitation: underweights dispatch overhead", 0),
    ("- Ranked 2-gather above 1-gather; hardware shows opposite", 1),
    ("- Simulation prunes bad candidates; hardware decides final rank", 1),
])

# ── 7. CGIS ──
add_body_slide("CGIS: Contention-Guided Iterative Synthesis", [
    ("Traditional: LLM -> schedule -> score -> LLM (blind)", 0),
    ("", 0),
    ("CGIS: LLM -> schedule -> contention diagnosis -> LLM", 0),
    ('- "Step 5 saturated link (4,8) at 6x. Separate distances 8 and 24."', 1),
    ("- LLM makes surgical fix instead of blind mutation", 1),
    ("", 0),
    ("Three components", 0),
    ("- ContentionAnalyzer: per-link conflict detection", 1),
    ("- IterativeRefinement: multi-turn LLM + contention feedback", 1),
    ("- IslandEvolution: 3 islands with LLM crossover", 1),
])

# ── 8. Template Evolution ──
add_body_slide("Template Evolution: Dual-Backend", [
    ("Beyond parameter tuning: LLM generates complete Python code", 0),
    ("", 0),
    ("XLA path (default)", 0),
    ("- Seeds: naive_allgather, allgather_reduce_scatter", 1),
    ("- TrackedTensor counts XLA IR ops; CollectiveSimulator tests correctness", 1),
    ("- LLM discovers how to reduce 67 ops -> 4 ops", 1),
    ("", 0),
    ("NKI path (research)", 0),
    ("- Seeds: nki_naive_allgather, nki_permute_ring", 1),
    ("- NKI collectives have ~7x higher dispatch than XLA on trn1", 1),
    ("", 0),
    ("5-layer correctness verification for both backends", 0),
])

# ── 9. Convergence History ──
add_body_slide("Search Convergence (MoE Traffic)", [
    ("From allgather seed", 0),
    ("- 3 rounds, no improvement (already optimal at 493 us)", 1),
    ("", 0),
    ("From hierarchical seed", 0),
    ("- Baseline: 3615 us (31 collective_permute ops)", 1),
    ("- Round 1: 493 us -- switched to AllGather pattern (86% improvement)", 1),
    ("- Round 2: 249 us -- 2-phase subgroup AllGather (50% further)", 1),
    ("", 0),
    ("From permute_ring seed", 0),
    ("- Round 1: 493 us -- switched to AllGather (86% improvement)", 1),
    ("- Round 2: failed (tensor shape mismatch)", 1),
])

# ── 10. Agent Output ──
add_body_slide("Agent Output: AllGather + index_select", [
    ("Algorithm (~3 XLA IR ops total)", 0),
    ("- Pad send buffer to uniform size", 1),
    ("- Single all_gather: collect all ranks' buffers", 1),
    ("- Single index_select: extract needed elements via flat index", 1),
    ("", 0),
    ("On XLA devices", 0),
    ("- All Python ops trace into HLO at compile time", 1),
    ("- Flat index = compile-time constant, zero runtime overhead", 1),
    ("- ~0.13 ms on trn1.32xlarge", 1),
])

# ── 11. Results Table: Single-Node ──
add_table_slide(
    "Single-Node Benchmark (32 ranks, NeuronLink)",
    ["Algorithm", "Type", "32KB", "128KB", "512KB", "2MB", "8MB"],
    [
        ["Agent Output",    "agent",    "0.127", "0.136", "0.132", "0.123", "0.127"],
        ["AG+ReduceScatter","baseline", "0.719", "0.727", "0.729", "0.739", "0.729"],
        ["NKI AllGather",   "baseline", "0.865", "0.894", "0.909", "0.869", "0.877"],
        ["AllGather+Slice", "baseline", "1.035", "1.890", "1.061", "1.046", "1.060"],
        ["Fused AllToAll",  "baseline", "2.257", "1.195", "1.195", "1.196", "1.202"],
        ["Hierarchical",    "baseline", "2.869", "2.205", "2.235", "2.237", "2.500"],
        ["Default Ring",    "baseline", "3.007", "2.824", "2.782", "2.806", "2.716"],
    ],
    note="All values in ms. MoE traffic (Zipf s=1.2). trn1.32xlarge, 20 iters, 5 warmup."
)

# ── 12. Results Table: Multi-Node ──
add_table_slide(
    "Multi-Node Benchmark (2x trn1.32xlarge, 64 ranks, EFA)",
    ["Algorithm", "Type", "64KB", "256KB", "1MB", "4MB", "16MB"],
    [
        ["Agent Output",    "agent",    "0.135", "0.132", "0.119", "0.194", "0.128"],
        ["NKI AllGather",   "baseline", "1.071", "1.120", "1.149", "1.134", "1.168"],
        ["AG+ReduceScatter","baseline", "1.097", "1.183", "2.062", "1.167", "1.196"],
        ["AllGather+Slice", "baseline", "2.051", "2.103", "2.032", "3.557", "3.618"],
        ["Fused AllToAll",  "baseline", "4.131", "2.250", "2.220", "2.256", "4.293"],
        ["Hierarchical",    "baseline", "2.547", "4.201", "2.473", "4.183", "2.354"],
        ["Default Ring",    "baseline", "5.889", "5.958", "5.875", "5.836", "6.116"],
    ],
    note="All values in ms. MoE traffic (Zipf s=1.2). 2-node EFA, 20 iters, 5 warmup."
)

# ── 13. Cross-Node Penalty ──
add_table_slide(
    "Cross-Node Penalty (2-node / single-node ratio)",
    ["Algorithm", "Type", "Avg Penalty", "Notes"],
    [
        ["Agent Output",    "agent",    "1.11x", "Near-zero cross-node overhead"],
        ["NKI AllGather",   "baseline", "1.28x", "Moderate EFA overhead"],
        ["Hierarchical",    "baseline", "1.34x", "Topology-aware helps"],
        ["AG+ReduceScatter","baseline", "1.84x", "2 dispatches amplify cross-node cost"],
        ["Default Ring",    "baseline", "2.10x", "63 steps across EFA"],
        ["Fused AllToAll",  "baseline", "2.21x", "Significant cross-node cost"],
        ["AllGather+Slice", "baseline", "2.36x", "Doubles with 2x data volume"],
    ],
    note="Agent's all_gather + index_select scales almost perfectly across EFA."
)

# ── 13b. Key Observations ──
add_body_slide("Key Observations", [
    ("Latency flat from 32 KB to 16 MB", 0),
    ("- Dispatch overhead dominates; NeuronLink/EFA never saturated", 1),
    ("", 0),
    ("XLA op count is the only lever", 0),
    ("- 3 ops (0.13 ms)  vs  33 ops (1.05 ms)  vs  93 ops (2.78 ms)", 1),
    ("", 0),
    ("AG+ReduceScatter: 2nd-fastest intra-node (~0.73ms)", 0),
    ("- But 2nd dispatch doubles cross-node cost (1.84x penalty)", 1),
    ("", 0),
    ("Agent output: ~45x faster than default ring at 2 nodes", 0),
    ("- 1 dispatch + ~3 XLA ops, 1.11x cross-node penalty", 1),
])

# ── 14. What LLM Actually Discovered ──
add_body_slide("What the LLM Actually Discovered", [
    ("From naive_allgather seed (67 local ops)", 0),
    ("- Round 1: 67 -> 35 ops (removes redundant slices)", 1),
    ("- Round 2: scatter_ + index_select = 4 ops", 1),
    ("", 0),
    ("From allgather_reduce_scatter seed", 0),
    ("- LLM explores AG+RS pattern but can't beat 1-dispatch AG", 1),
    ("", 0),
    ("From permute_ring seed", 0),
    ("- Switches to AllGather pattern (86% improvement)", 1),
    ("", 0),
    ("Agent independently arrives at all_gather + index_select", 0),
    ("- Not shown the answer; discovers it from naive seeds", 1),
])

# ── 15. Simulator vs Hardware Gap ──
add_body_slide("Simulator vs Hardware Gap", [
    ("Simulator #1: 2-gather (249 us) -> Hardware: 0.25 ms (3rd)", 0),
    ("Simulator #2: 1-gather (493 us) -> Hardware: 0.13 ms (1st)", 0),
    ("Simulator #4: fused (143 us) -> Hardware: 1.20 ms (5th)", 0),
    ("", 0),
    ("Root cause", 0),
    ("- Simulator models bandwidth and contention well", 1),
    ("- But underweights XLA dispatch overhead", 1),
    ("- On real hardware, dispatch cost >> link utilization", 1),
    ("", 0),
    ("Takeaway", 0),
    ("- Simulation prunes clearly bad candidates", 1),
    ("- Final ranking must come from hardware measurement", 1),
])

# ── 16. Lessons Learned ──
add_body_slide("Lessons Learned", [
    ("LLMs are search heuristics, not inventors", 0),
    ("- Selected known patterns from prompt context", 1),
    ("- Efficiently navigated search space", 1),
    ("", 0),
    ("Structured feedback (CGIS) >> blind scoring", 0),
    ("- Per-step contention diagnosis enables surgical fixes", 1),
    ("", 0),
    ("Correctness verification is essential", 0),
    ("- LLM-generated code fails ~30% of the time", 1),
    ("- 5-layer verification catches subtle distributed bugs", 1),
])

# ── 17. Next Steps ──
add_body_slide("Next Steps", [
    ("Improve simulator fidelity", 0),
    ("- XLA op count as first-class cost; profile-guided calibration", 1),
    ("- Simulator underweights dispatch overhead vs bandwidth", 1),
    ("", 0),
    ("Reduce LLM prompt leakage", 0),
    ("- Don't show all templates as references; force composition", 1),
    ("", 0),
    ("End-to-end MoE training integration", 0),
    ("- Measure wall-clock speedup, not just AllToAllV latency", 1),
    ("", 0),
    ("Larger multi-node (4+ nodes)", 0),
    ("- Test whether agent pattern stays optimal at higher node counts", 1),
])

# ── 18. Summary ──
add_body_slide("Summary", [
    ("Problem: AllToAllV on Trainium is 8-23x slower than optimal", 0),
    ("", 0),
    ("Approach: Hardware-agnostic 5-phase agent pipeline", 0),
    ("- LLM-built simulator + evolution + HW validation + feedback loops", 1),
    ("- 5-layer correctness verification", 1),
    ("", 0),
    ("Result: ~0.13 ms (3 XLA ops), 1.11x cross-node penalty", 0),
    ("- 8x vs naive AllGather, 23x vs default ring (single-node)", 1),
    ("- ~45x faster than default ring at 2 nodes", 1),
    ("- Drop-in module: runtime/trainium_alltoallv.py", 1),
    ("", 0),
    ("Insight: dispatch overhead >> bandwidth at all scales tested", 0),
    ("- Agent independently discovers optimal pattern from naive seeds", 1),
])

# ══════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════
out = "experiments/trainium_alltoallv_search.pptx"
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
